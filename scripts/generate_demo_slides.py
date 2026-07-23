"""Generate AIVIA Demo Slide Deck v2.

Includes ecosystem positioning and GPS analogy.

Run: python3 scripts/generate_demo_slides.py
Output: presentation/aivia_demo.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Color palette
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DEEP_NAVY = RGBColor(0x0D, 0x1B, 0x2A)
TEAL = RGBColor(0x00, 0x96, 0x88)
LIGHT_TEAL = RGBColor(0x4D, 0xB6, 0xAC)
ORANGE = RGBColor(0xFF, 0x7A, 0x2F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MEDIUM_GRAY = RGBColor(0x75, 0x75, 0x75)
DARK_RED = RGBColor(0xC6, 0x28, 0x28)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
DARK_BG = RGBColor(0x2A, 0x1A, 0x1A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = Inches(13.333)


def bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def txt(slide, l, t, w, h, s, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = s
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = "Calibri"
    p.alignment = a


def bullets(slide, l, t, w, h, items, sz=18, c=WHITE, bc=TEAL, sp=8):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(sp)
        rb = p.add_run()
        rb.text = "\u25B8 "
        rb.font.size = Pt(sz)
        rb.font.color.rgb = bc
        rb.font.name = "Calibri"
        if isinstance(item, tuple):
            r1 = p.add_run()
            r1.text = item[0]
            r1.font.size = Pt(sz)
            r1.font.color.rgb = c
            r1.font.bold = True
            r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = item[1]
            r2.font.size = Pt(sz)
            r2.font.color.rgb = c
            r2.font.name = "Calibri"
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(sz)
            r.font.color.rgb = c
            r.font.name = "Calibri"


def rrect(slide, l, t, w, h, fc=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if fc:
        s.fill.solid()
        s.fill.fore_color.rgb = fc
    return s


def rect(slide, l, t, w, h, fc=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if fc:
        s.fill.solid()
        s.fill.fore_color.rgb = fc
    return s


def arrow_r(slide, l, t, w, h, fc=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fc
    s.line.fill.background()


def arrow_d(slide, l, t, w, h, fc=ORANGE):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fc
    s.line.fill.background()


# ============================================================
# SLIDE 1 — Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)

circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-1), Inches(7), Inches(7))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x22, 0x3A, 0x5E)
circle.line.fill.background()

rect(s, Inches(0.8), Inches(4.4), Inches(4), Pt(4), TEAL)

txt(s, Inches(0.8), Inches(1.5), Inches(8), Inches(0.6),
    "AIVIA", sz=24, c=LIGHT_TEAL)
txt(s, Inches(0.8), Inches(2.2), Inches(10), Inches(1.5),
    "Turn Technical Debt\ninto Self-Service Intelligence", sz=48, c=WHITE, b=True)
txt(s, Inches(0.8), Inches(4.7), Inches(8), Inches(0.5),
    "The GPS Map for Microsoft Fabric's Conversational Analytics", sz=20, c=LIGHT_TEAL)
txt(s, Inches(0.8), Inches(5.8), Inches(8), Inches(0.3),
    "www.aiviaapp.com", sz=18, c=TEAL)


# ============================================================
# SLIDE 2 — The Problem (Hidden Treasure)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)

txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "The Hidden Treasure Problem", sz=40, c=WHITE, b=True)
rect(s, Inches(0.8), Inches(1.3), Inches(4), Pt(4), TEAL)

# Left — the treasure
txt(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.5),
    "What your organization already has:", sz=20, c=TEAL, b=True)
bullets(s, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.5), [
    ("Thousands", " of SQL-based reports"),
    ("Requested by clinicians", " based on real needs"),
    ("Built by skilled BI developers", ""),
    ("Validated", " and in production for years"),
    ("= Millions of dollars", " of invested business logic"),
], sz=19)

# Right — the problem
rrect(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.8), DARK_BG)
txt(s, Inches(7.3), Inches(2.0), Inches(5), Inches(0.5),
    "But nobody can find or trust it", sz=22, c=DARK_RED, b=True)
bullets(s, Inches(7.3), Inches(2.7), Inches(5), Inches(3.8), [
    ("No governance", " — no visibility into the logic"),
    ("Outdated reports", " next to current ones"),
    ("Conflicting duplicates", " — same metric, different numbers"),
    ("Clinicians staring", " at reports they can't trust"),
    ("Submit another request.", " Wait weeks. More distrust."),
], sz=17, bc=DARK_RED)

# Bottom
rect(s, Inches(0), Inches(6.5), SW, Inches(1), TEAL)
txt(s, Inches(1), Inches(6.6), Inches(11), Inches(0.6),
    "This is the cycle we break.", sz=28, c=WHITE, b=True, a=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3 — The GPS Analogy
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)

txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Where AIVIA Fits in Microsoft Fabric", sz=36, c=WHITE, b=True)
rect(s, Inches(0.8), Inches(1.3), Inches(4), Pt(4), TEAL)

# Microsoft's stack
stack_items = [
    ("Conversational Analytics", "Data Agents, Copilot", Inches(1.8), ORANGE),
    ("Semantic Layer", "Power BI Models", Inches(3.0), RGBColor(0x00, 0x79, 0x6B)),
    ("Serving Layer", "Data Warehouse", Inches(4.2), RGBColor(0x26, 0xA6, 0x9A)),
    ("Data Foundation", "Lakehouse, Pipelines", Inches(5.4), TEAL),
]

for label, desc, top, color in stack_items:
    rrect(s, Inches(0.8), top, Inches(5), Inches(1.0), color)
    txt(s, Inches(1.0), top + Inches(0.1), Inches(2.5), Inches(0.4),
        label, sz=16, c=WHITE, b=True)
    txt(s, Inches(1.0), top + Inches(0.45), Inches(2.5), Inches(0.4),
        desc, sz=13, c=WHITE)

# AIVIA arrow pointing between layers 3 and 4
rrect(s, Inches(6.2), Inches(2.2), Inches(6.3), Inches(2.5), RGBColor(0x22, 0x3A, 0x5E))
txt(s, Inches(6.5), Inches(2.3), Inches(5.7), Inches(0.5),
    "AIVIA = The GPS Map", sz=26, c=ORANGE, b=True)
txt(s, Inches(6.5), Inches(2.9), Inches(5.7), Inches(1.5),
    "Microsoft built the highway (Fabric)\n"
    "Microsoft built the car (Data Agent / Copilot)\n\n"
    "AIVIA builds the map — the knowledge of where\n"
    "everything is, what it means, and how to get there.\n\n"
    "Without the map, the car drives blind.",
    sz=16, c=WHITE)

# Bottom note
txt(s, Inches(6.5), Inches(5.2), Inches(5.7), Inches(1.5),
    "Microsoft provides empty containers.\n"
    "AIVIA fills them with certified, traceable business logic.",
    sz=16, c=LIGHT_TEAL)


# ============================================================
# SLIDE 4 — How It Works (Pipeline)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)

txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "How AIVIA Works", sz=40, c=WHITE, b=True)
rect(s, Inches(0.8), Inches(1.3), Inches(4), Pt(4), TEAL)

txt(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
    "Take your existing SQL and automatically build a certified knowledge graph",
    sz=20, c=LIGHT_TEAL)

pipeline = [
    ("Your SQL", "Stored Procs\nViews\nScripts", Inches(0.3), MEDIUM_GRAY),
    ("ScriptDom", "Microsoft's own\nT-SQL parser\n99% accuracy", Inches(2.8), TEAL),
    ("Knowledge\nGraph", "3-Layer certified\nbusiness logic map", Inches(5.3), RGBColor(0x26, 0xA6, 0x9A)),
    ("Data Agent", "Ask questions\nin plain English", Inches(7.8), RGBColor(0x00, 0x79, 0x6B)),
    ("Purview\nCollibra\nPower BI", "Auto-governed\nAlways current", Inches(10.3), ORANGE),
]

for label, desc, left, color in pipeline:
    rrect(s, left, Inches(2.8), Inches(2.2), Inches(2.5), color)
    txt(s, left + Inches(0.1), Inches(2.9), Inches(2.0), Inches(0.8),
        label, sz=16, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(s, left + Inches(0.1), Inches(3.7), Inches(2.0), Inches(1.2),
        desc, sz=12, c=WHITE, a=PP_ALIGN.CENTER)

for x in [Inches(2.5), Inches(5.0), Inches(7.5), Inches(10.0)]:
    arrow_r(s, x, Inches(3.8), Inches(0.4), Inches(0.3))

# Stats bar
rect(s, Inches(0), Inches(5.8), SW, Inches(1.5), DEEP_NAVY)
for val, label, left, color in [
    ("1,300+", "SQL sources parsed", Inches(1), ORANGE),
    ("99%", "parse accuracy", Inches(4.8), GREEN),
    ("0", "parse errors", Inches(8.6), TEAL),
]:
    txt(s, left, Inches(5.9), Inches(3.5), Inches(1.0),
        val, sz=54, c=color, b=True, a=PP_ALIGN.CENTER)
    txt(s, left, Inches(6.7), Inches(3.5), Inches(0.4),
        label, sz=18, c=MEDIUM_GRAY, a=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5 — Live Demo
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)

txt(s, Inches(1), Inches(2.5), Inches(11), Inches(1.0),
    "Live Demo", sz=48, c=WHITE, b=True, a=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
    "Ask the Data Agent anything about your metrics — in plain English",
    sz=24, c=LIGHT_TEAL, a=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
    "[Switch to Fabric Data Agent screen recording]",
    sz=18, c=MEDIUM_GRAY, a=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6 — Knowledge Graph Visual
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)

txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "The Knowledge Graph Behind Every Answer", sz=36, c=WHITE, b=True)
rect(s, Inches(0.8), Inches(1.3), Inches(4), Pt(4), TEAL)

layers = [
    ("BUSINESS METRICS", "Certified definitions: ER Length of Stay, Census Dashboard...",
     Inches(2.0), TEAL),
    ("CALCULATION LOGIC", "SQL fragments: DATEDIFF, COUNT, AVG, WHERE filters...",
     Inches(3.5), RGBColor(0x26, 0xA6, 0x9A)),
    ("SOURCE TABLES", "Physical tables: CLARITY_ADT, PATIENT, PAT_ENC...",
     Inches(5.0), RGBColor(0x00, 0x60, 0x56)),
]

for label, desc, top, color in layers:
    rrect(s, Inches(2), top, Inches(9), Inches(1.2), color)
    txt(s, Inches(2.3), top + Inches(0.1), Inches(8.4), Inches(0.4),
        label, sz=18, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(s, Inches(2.3), top + Inches(0.5), Inches(8.4), Inches(0.5),
        desc, sz=14, c=WHITE, a=PP_ALIGN.CENTER)

for y in [Inches(3.2), Inches(4.7)]:
    arrow_d(s, Inches(6.3), y, Inches(0.4), Inches(0.3))

txt(s, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
    "Every answer traces from business question \u2192 calculation logic \u2192 source table",
    sz=18, c=MEDIUM_GRAY, a=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7 — Why Not Just Use Copilot?
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)

txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Why Not Just Use Copilot?", sz=36, c=WHITE, b=True)
rect(s, Inches(0.8), Inches(1.3), Inches(4), Pt(4), TEAL)

# Left — what Microsoft provides
rrect(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), RGBColor(0x22, 0x3A, 0x5E))
txt(s, Inches(1.0), Inches(2.0), Inches(5.1), Inches(0.5),
    "What Microsoft Provides", sz=22, c=LIGHT_TEAL, b=True)
bullets(s, Inches(1.0), Inches(2.7), Inches(5.1), Inches(3.5), [
    ("Data Agent", " — empty chat interface"),
    ("Purview", " — empty catalog"),
    ("Power BI", " — empty report descriptions"),
    ("Copilot", " — generic AI, no org knowledge"),
    ("ScriptDom", " — parser library (just a tool)"),
], sz=17, bc=LIGHT_TEAL)

# Right — what AIVIA adds
rrect(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5), RGBColor(0x1A, 0x2A, 0x1A))
txt(s, Inches(7.2), Inches(2.0), Inches(5.1), Inches(0.5),
    "What AIVIA Fills In", sz=22, c=GREEN, b=True)
bullets(s, Inches(7.2), Inches(2.7), Inches(5.1), Inches(3.5), [
    ("Knowledge graph", " grounding the agent"),
    ("Business term definitions", " pushed to Purview"),
    ("Auto-generated descriptions", " for every report"),
    ("Org-specific certified logic", " — not generic AI"),
    ("99% accurate parsing", " — ScriptDom + intelligence"),
], sz=17, bc=GREEN)

# Bottom
rect(s, Inches(0), Inches(6.5), SW, Inches(1), DEEP_NAVY)
txt(s, Inches(1), Inches(6.6), Inches(11), Inches(0.6),
    "Microsoft builds empty containers. AIVIA fills them with certified business logic.",
    sz=22, c=ORANGE, b=True, a=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 8 — Roadmap & Close
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)

txt(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "What's Next", sz=40, c=WHITE, b=True)
rect(s, Inches(0.8), Inches(1.3), Inches(4), Pt(4), TEAL)

roadmap = [
    ("1", "Metadata Sync", "Auto-populate Purview & Power BI descriptions",
     TEAL),
    ("2", "Governance Flywheel", "Every question strengthens the knowledge base",
     RGBColor(0x26, 0xA6, 0x9A)),
    ("3", "Multi-Dialect", "Oracle PL/SQL, Snowflake — native parsers per platform",
     RGBColor(0x00, 0x60, 0x56)),
]

for i, (num, title, desc, color) in enumerate(roadmap):
    top = Inches(2.0) + Inches(1.5) * i
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), top + Inches(0.1),
                              Inches(0.6), Inches(0.6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    txt(s, Inches(1.0), top + Inches(0.12), Inches(0.6), Inches(0.6),
        num, sz=22, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(s, Inches(1.9), top, Inches(10), Inches(0.5),
        title, sz=24, c=WHITE, b=True)
    txt(s, Inches(1.9), top + Inches(0.5), Inches(10), Inches(0.5),
        desc, sz=17, c=LIGHT_TEAL)

# Close
rect(s, Inches(0), Inches(6.0), SW, Inches(1.5), TEAL)
txt(s, Inches(1), Inches(6.1), Inches(11), Inches(0.5),
    "Entirely inside your Microsoft Fabric tenant. No data ever leaves.",
    sz=20, c=WHITE, a=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.6), Inches(11), Inches(0.6),
    "AIVIA is live, scalable, and ready for deployment.",
    sz=22, c=WHITE, b=True, a=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(7.1), Inches(11), Inches(0.3),
    "www.aiviaapp.com", sz=18, c=NAVY, b=True, a=PP_ALIGN.CENTER)


# ============================================================
# Save
# ============================================================
output_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "presentation")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "aivia_demo.pptx")
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
