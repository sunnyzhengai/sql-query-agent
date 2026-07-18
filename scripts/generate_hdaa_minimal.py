"""
Generate HDAA deck with minimal formatting — plain text on blank slides.
Easy to copy/paste into the HDAA template.

Run: python3 scripts/generate_hdaa_minimal.py
Output: presentation/hdaa_full_circle_minimal.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_text(slide, left, top, width, height, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=BLACK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        if isinstance(item, tuple):
            run = p.add_run()
            run.text = item[0]
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = True
            run2 = p.add_run()
            run2.text = item[1]
            run2.font.size = Pt(size)
            run2.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = "- " + item
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


# Speaker notes text for each slide
notes = [
    # Slide 1
    "Welcome everyone. I'm Sunny Zheng, and today I want to talk about the tension between governance and getting work done. The title is 'The Full Circle' because what I'm going to show you is a system where the more people use it, the better it gets. And the governance happens automatically.",
    # Slide 2
    "I lead BI in healthcare — SQL, data warehousing, making data trustworthy. We work in Microsoft Fabric. The challenge: over a thousand reports, business logic scattered everywhere, users waiting weeks for answers. Every time we try to govern our way out of it, we slow everyone down.",
    # Slide 3
    "Here's our roadmap. Problem, big idea, foundation, flywheel, a real case study with surgeons, ROI, and lessons learned.",
    # Slide 4
    "Sound familiar? Raise your hand if you've heard 'can you pull me the readmission rate' and it's a 6-week wait. Three analysts, three answers. Governance is supposed to fix this but it's a committee that meets monthly. The vicious cycle: governance is slow, people bypass it, definitions go stale, trust declines, more governance imposed. Repeat.",
    # Slide 5
    "What if every question a user asks makes your data foundation stronger? Instead of governance being a gate, it becomes a knowledge base that grows from demand. Every question reinforces what we know or surfaces what we don't. Governance becomes invisible to the end user.",
    # Slide 6
    "Before the flywheel can spin, you need a starting point. The good news? You already have the raw material.",
    # Slide 7
    "Your existing queries already contain the business logic — it's just buried in SQL. We scan automatically, AI drafts definitions, team reviews in bulk, stewards certify. One OR dashboard yields three candidate metrics, all linked to source. This is how you go from zero to 50+ certified metrics without committee meetings.",
    # Slide 8
    "Three layers make answers traceable. Top: business metrics, owned by stewards, weighted by usage. Middle: the calculation logic — small reusable pieces, not giant SQL blocks. Bottom: source data with plain-English descriptions. Every answer traces from top to bottom. No black boxes.",
    # Slide 9
    "Now here's where it gets interesting.",
    # Slide 10
    "The flywheel: user asks, agent answers, knowledge base grows, better answers, more questions. Notice what's NOT here: committees, spreadsheets, project proposals. It grows because people use it.",
    # Slide 11
    "Two paths, both good. Path A: we know this — instant answer, metric gains weight, most-asked metrics become visible. Path B: we don't know yet — agent says 'I don't have that yet,' steward gets notified, certifies it, next time it's Path A. No wasted questions.",
    # Slide 12
    "Let me show you a real example.",
    # Slide 13
    "FCOTS: percentage of first surgical cases starting on time. Three years ago: 20%. Today: 60%+. Goal: 90%+. Working with the Surgery Medical Director. This improvement was only possible because everyone trusted the same number. One certified definition.",
    # Slide 14
    "The Medical Director needs 1:1s with surgeons about their performance. But he can't share a dashboard showing everyone's numbers — it's sensitive. Our solution: surgeons ask the agent 'What's my FCOTS?' and only see their own data. Dr. Smith sees 78%, Dr. Jones sees 62%, the Medical Director sees the department view. Same metric, personalized access, no IT tickets, no privacy concerns.",
    # Slide 15
    "Watch the flywheel spin. Surgeons ask about FCOTS — Path A. Then they ask new questions: 'What's my delay time?' 'Which rooms start late?' Path B — steward certifies new metrics. One metric grows into a full OR efficiency suite, driven by surgeon demand. Not a governance project — organic growth from clinical needs.",
    # Slide 16
    "The FCOTS example is one department. Imagine this across your entire organization.",
    # Slide 17
    "Month 1: seed 50 metrics from existing reports. Month 3: usage adds 30 more. Month 6: 140. Month 12: 250+ and still growing. Traditional governance: 50 per year if you're lucky. The flywheel: 250+ and accelerating.",
    # Slide 18
    "Network effects. More metrics, more answers, more trust, more users, more questions, more certifications. Week 1: most questions are new. Month 3: 80% instant. Month 12: 95%. Steward workload actually decreases over time as coverage grows.",
    # Slide 19
    "How do users experience this? Two channels working together.",
    # Slide 20
    "The agent handles ad-hoc: no ticket, no wait, personalized security. Dashboards handle recurring KPIs. Both use the same certified knowledge base — numbers always match. The agent's usage data tells you which dashboards to build. When FCOTS gets asked 100 times, that's your signal.",
    # Slide 21
    "Let's talk numbers.",
    # Slide 22
    "Four metrics: knowledge base coverage 50 to 250+. Instant answer rate 40% to 95%. Time to answer: weeks to seconds. Steward efficiency: 10x. All compound over time.",
    # Slide 23
    "Before: which LOS is right? Three answers. Six-week backlog. Committee governance. After: certified ER LOS asked 347 times. One definition. Seconds to answer. Governance happens automatically. Usage shows what the org cares about.",
    # Slide 24
    "Six lessons. Start with what you have. 'I don't know' is a feature. Let usage drive prioritization. Make governance invisible. Security enables trust. The flywheel compounds.",
    # Slide 25
    "Four takeaways. Your reports are an untapped asset. Usage is the best governance signal. Self-service and governance can be the same thing. The full circle compounds over time.",
    # Slide 26
    "Thank you. I'd love your questions — especially if you're dealing with similar challenges.",
]


# -- SLIDE 1: Title --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(0.8),
         "The Full Circle:", size=28, color=GRAY)
add_text(slide, Inches(0.8), Inches(2.3), Inches(10), Inches(1.5),
         "From SQL Chaos to Self-Service AI Agent", size=44, bold=True)
add_text(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.5),
         "How Every User Question Builds Your Certified Data Foundation", size=20, color=GRAY)
add_text(slide, Inches(0.8), Inches(5.5), Inches(10), Inches(0.4),
         "Sunny Zheng", size=22, bold=True)
add_text(slide, Inches(0.8), Inches(6.0), Inches(10), Inches(0.4),
         "HDAA Annual Conference  |  November 2026", size=16, color=GRAY)

# -- SLIDE 2: About --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "About the Speaker", size=30, bold=True)
add_bullets(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4), [
    ("BI Management Leader", " in Healthcare"),
    ("15+ years", " turning data into decisions for clinical and operational teams"),
    ("Microsoft Fabric ecosystem", " — Power BI, Data Agent, Notebooks"),
    ("Focused on", " making analytics self-service without sacrificing accuracy"),
    ("Passionate about", " making governance disappear into the background"),
], size=18)
add_text(slide, Inches(7.5), Inches(1.5), Inches(4.5), Inches(0.5),
         "The Challenge", size=22, bold=True)
add_bullets(slide, Inches(7.5), Inches(2.2), Inches(4.5), Inches(3), [
    "1,000+ reports and queries",
    "Business logic buried in code",
    "Users waiting weeks for answers",
    "Governance slowing everyone down",
], size=18)

# -- SLIDE 3: Agenda --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Agenda", size=30, bold=True)
add_bullets(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5), [
    "1. The Problem: Why Governance Feels Like a Tax",
    "2. The Idea: What If Usage Drove Governance?",
    "3. Building the Foundation: From Existing SQL to a Certified Knowledge Base",
    "4. The Flywheel: Ask, Answer, Certify, Grow",
    "5. Case Study: First Case On Time Start (FCOTS)",
    "6. The Growing Knowledge Base: Governance as a Byproduct",
    "7. Delivering Value: Self-Service + Dashboards",
    "8. Measuring the Impact: ROI That Compounds",
    "9. Lessons Learned & Q&A",
], size=20)

# -- SLIDE 4: The Problem --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "The Problem: Why Governance Feels Like a Tax", size=30, bold=True)
add_text(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
         "Sound familiar?", size=20, bold=True, color=GRAY)
add_bullets(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.5), [
    "\"Can you pull me the readmission rate?\" — 6-week backlog",
    "Three analysts run the same metric, get three answers",
    "Governance team creates definitions nobody uses",
    "The more you govern, the slower everyone moves",
    "So governance falls behind, definitions go stale, trust erodes",
    "Result: governance and getting work done feel like opposites",
], size=17)
add_text(slide, Inches(7.5), Inches(1.3), Inches(4.5), Inches(0.5),
         "The Vicious Cycle", size=22, bold=True)
add_bullets(slide, Inches(7.5), Inches(2.0), Inches(4.5), Inches(3.5), [
    "Governance is slow",
    "People work around it",
    "Definitions go stale",
    "Trust in data declines",
    "More governance imposed",
    "(repeat forever)",
], size=17)

# -- SLIDE 5: The Big Idea --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "The Idea: What If Usage Drove Governance?", size=30, bold=True)
add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
         "What if every question a user asks\nmakes your data foundation stronger?",
         size=32, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(2.5),
         "Instead of governance being a gate people pass through,\n"
         "it becomes a knowledge base that grows from demand.\n\n"
         "Every question either reinforces what we know\n"
         "or surfaces what we don't — and kicks off a review.",
         size=20, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.5),
         "Governance is no longer a chore. It's a byproduct of people doing their jobs.",
         size=18, bold=True, align=PP_ALIGN.CENTER)

# -- SLIDE 6: Section Divider --
slide = blank_slide()
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "BUILDING THE FOUNDATION", size=40, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "From Existing SQL to a Certified Knowledge Base", size=22, color=GRAY, align=PP_ALIGN.CENTER)

# -- SLIDE 7: Foundation --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "The Foundation: Your SQL Already Contains the Answers", size=30, bold=True)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
         "You don't start from scratch — your existing queries are the raw material", size=18, color=GRAY)
add_bullets(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), [
    ("Scan your existing reports and queries", " — automatically extract the business logic"),
    ("AI drafts plain-English definitions", " from the logic it finds"),
    ("Your team reviews in bulk", " — not one definition at a time"),
    ("Stewards certify", " the business meaning: approve, modify, or reject"),
    ("Result:", " an initial certified knowledge base of your most-used metrics"),
    ("This is the seed.", " What makes it grow is what happens next."),
], size=17)
add_text(slide, Inches(7.5), Inches(2.0), Inches(4.5), Inches(0.4),
         "Example: One report yields...", size=16, color=GRAY)
add_text(slide, Inches(7.5), Inches(2.6), Inches(4.5), Inches(3.5),
         "Report: OR Efficiency Dashboard\n\n"
         "  3 business metrics found:\n"
         "    \"First Case On Time Start\"\n"
         "    \"Turnover Time\"\n"
         "    \"Case Duration Accuracy\"\n\n"
         "  Each linked to its source\n"
         "  tables and calculation logic\n\n"
         "  -> Sent to stewards for review",
         size=14, color=GRAY)

# -- SLIDE 8: Knowledge Base Structure --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "The Foundation: How the Knowledge Base Is Organized", size=30, bold=True)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
         "Three layers connect business questions to source data — with full traceability", size=18, color=GRAY)
add_text(slide, Inches(2), Inches(2.2), Inches(9), Inches(0.4),
         "BUSINESS METRICS", size=20, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(2.7), Inches(9), Inches(0.5),
         "Certified definitions (e.g., First Case On Time Start)\nOwned by a steward  |  Weighted by how often people ask",
         size=15, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(6.3), Inches(3.3), Inches(0.5), Inches(0.4),
         "v", size=24, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.4),
         "CALCULATION LOGIC", size=20, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(4.3), Inches(9), Inches(0.5),
         "The step-by-step logic behind each metric\nSmall, reusable pieces — not giant SQL blocks",
         size=15, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(6.3), Inches(4.9), Inches(0.5), Inches(0.4),
         "v", size=24, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(5.4), Inches(9), Inches(0.4),
         "SOURCE DATA", size=20, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(5.9), Inches(9), Inches(0.5),
         "The actual tables, columns, and filters\nEnriched with plain-English descriptions",
         size=15, color=GRAY, align=PP_ALIGN.CENTER)

# -- SLIDE 9: Section Divider --
slide = blank_slide()
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "THE FLYWHEEL", size=40, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "Ask, Answer, Certify, Grow", size=22, color=GRAY, align=PP_ALIGN.CENTER)

# -- SLIDE 10: Flywheel --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "The Flywheel: Every Question Makes the System Smarter", size=30, bold=True)
add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(3),
         "USER ASKS A QUESTION\n"
         "        |\n"
         "        v\n"
         "AGENT ANSWERS  ------>  KNOWLEDGE BASE GROWS\n"
         "        ^                       |\n"
         "        |                       v\n"
         "BETTER ANSWERS  <------  MORE CERTIFIED METRICS",
         size=18, align=PP_ALIGN.CENTER)
add_text(slide, Inches(4), Inches(3.2), Inches(5), Inches(0.6),
         "FULL CIRCLE", size=28, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
         "No governance committees. No spreadsheets. The knowledge base grows because people use it.",
         size=16, color=GRAY, align=PP_ALIGN.CENTER)

# -- SLIDE 11: Two Paths --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Two Paths: Every Question Has a Destination", size=30, bold=True)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
         "When someone asks a question, one of two things happens:", size=18, color=GRAY)
add_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.5),
         "PATH A: We Already Know This", size=22, bold=True)
add_bullets(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(3.5), [
    "The answer exists in the certified knowledge base",
    "User gets an instant, accurate answer",
    "The metric gains weight — it's clearly important",
    "Most-asked metrics become visible to leadership",
    "High-demand metrics get promoted to dashboards",
], size=16)
add_text(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(0.5),
         "PATH B: We Don't Know This Yet", size=22, bold=True)
add_bullets(slide, Inches(7), Inches(2.7), Inches(5.5), Inches(3.5), [
    "The agent says \"I don't have that yet\"",
    "A request is sent to the data steward",
    "Steward reviews and certifies the definition",
    "New metric added to the knowledge base",
    "Next time anyone asks, they get an instant answer",
], size=16)
add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         "Both paths make the system better. There is no wasted question.",
         size=18, bold=True, align=PP_ALIGN.CENTER)

# -- SLIDE 12: Section Divider --
slide = blank_slide()
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "CASE STUDY", size=40, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "First Case On Time Start (FCOTS)", size=22, color=GRAY, align=PP_ALIGN.CENTER)

# -- SLIDE 13: FCOTS Story --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Case Study: First Case On Time Start", size=30, bold=True)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
         "How one certified metric helped transform surgical efficiency", size=18, color=GRAY)
add_bullets(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), [
    ("FCOTS =", " % of first surgical cases that start on time"),
    ("3 years ago:", " our rate was just over 20%"),
    ("Today:", " we've reached 60%+"),
    ("Goal:", " 90%+ on-time start rate"),
    ("Working with", " the Surgery Medical Director to drive improvement"),
    ("One certified definition,", " one source of truth, consistent tracking over time"),
], size=17)
add_text(slide, Inches(7.5), Inches(2.0), Inches(4.5), Inches(0.5),
         "FCOTS Journey", size=22, bold=True)
add_text(slide, Inches(7.5), Inches(2.7), Inches(4.5), Inches(3),
         "3 Years Ago:  20%+\n\n"
         "Today:        60%+\n\n"
         "Goal:         90%+\n\n\n"
         "This improvement was only possible\n"
         "because everyone trusted the same number.",
         size=17, color=GRAY)

# -- SLIDE 14: FCOTS Self-Service --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "FCOTS: Why This Metric Needs Self-Service", size=30, bold=True)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
         "The Medical Director's challenge: individual performance is sensitive data", size=18, color=GRAY)
add_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.4),
         "The Challenge", size=20, bold=True)
add_bullets(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(1.5), [
    "Medical Director needs 1:1s with surgeons about their FCOTS",
    "Can't share a dashboard showing everyone's performance",
    "Individual surgeon data is sensitive — requires privacy",
], size=16)
add_text(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(0.4),
         "The Self-Service Solution", size=20, bold=True)
add_bullets(slide, Inches(0.8), Inches(4.9), Inches(5.5), Inches(1.5), [
    "Surgeons ask the Data Agent: \"What's my FCOTS rate?\"",
    "Each surgeon only sees their own data — built-in security",
    "Medical Director sees the department-level view",
], size=16)
add_text(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(0.5),
         "Same Metric, Right Access", size=20, bold=True)
add_text(slide, Inches(7), Inches(2.7), Inches(5.5), Inches(4),
         "Dr. Smith asks:\n"
         "  \"What's my FCOTS this month?\"\n"
         "  Answer: 78% (your cases only)\n\n"
         "Dr. Jones asks:\n"
         "  \"What's my FCOTS this month?\"\n"
         "  Answer: 62% (your cases only)\n\n"
         "Medical Director asks:\n"
         "  \"Show me department FCOTS\"\n"
         "  Answer: 68% overall + trends",
         size=15, color=GRAY)
add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
         "One certified definition. Personalized access. No IT tickets. No privacy concerns.",
         size=17, bold=True, align=PP_ALIGN.CENTER)

# -- SLIDE 15: FCOTS Flywheel --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "FCOTS: The Full Circle in Action", size=30, bold=True)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
         "What happens when surgeons start using the agent regularly:", size=18, color=GRAY)
add_bullets(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), [
    ("Surgeons ask about their FCOTS", " — the metric gains weight in the system"),
    ("Some ask new questions:", " \"What's my average delay time?\" \"Which rooms start late most?\""),
    ("Those questions trigger", " steward review — new metrics get certified"),
    ("The knowledge base grows", " around OR efficiency — driven by surgeon demand"),
    ("Medical Director sees", " which metrics surgeons care about most"),
    ("Result:", " a richer, more useful system built from actual clinical needs"),
], size=17)
add_text(slide, Inches(7.5), Inches(2.0), Inches(4.5), Inches(0.5),
         "FCOTS Flywheel Growth", size=20, bold=True)
add_text(slide, Inches(7.5), Inches(2.7), Inches(4.5), Inches(3.5),
         "Start:   1 certified metric (FCOTS)\n\n"
         "Month 1: Surgeons start asking\n"
         "         about delay reasons\n\n"
         "Month 3: 5 related OR efficiency\n"
         "         metrics certified\n\n"
         "Month 6: Full OR efficiency suite\n"
         "         built from surgeon demand",
         size=16, color=GRAY)

# -- SLIDE 16: Section Divider --
slide = blank_slide()
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "THE GROWING KNOWLEDGE BASE", size=40, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "Governance as a Byproduct of Work", size=22, color=GRAY, align=PP_ALIGN.CENTER)

# -- SLIDE 17: Growth --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "The Knowledge Base Grows From Demand, Not Committees", size=30, bold=True)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
         "Month over month, your certified metrics expand:", size=18, color=GRAY)
add_text(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(3),
         "Month 1:   50    Seed from existing reports & queries\n\n"
         "Month 3:   80    User questions add 30 new certified metrics\n\n"
         "Month 6:  140    Steward queue drives 60 more certifications\n\n"
         "Month 12: 250+   Flywheel matures",
         size=20)
add_text(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
         "Traditional governance: ~50 definitions/year.  Demand-driven: 250+ and accelerating.",
         size=18, bold=True, align=PP_ALIGN.CENTER)

# -- SLIDE 18: Acceleration --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Why the Flywheel Accelerates Over Time", size=30, bold=True)
add_bullets(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.5), [
    ("More certified metrics", " = more questions the agent can answer"),
    ("More answers", " = more people trust and adopt the agent"),
    ("More users", " = more questions — both known and new"),
    ("More new questions", " = more steward certifications"),
    ("More certifications", " = even more metrics the agent knows"),
    ("Each cycle is faster.", " The foundation compounds."),
], size=17)
add_text(slide, Inches(7.5), Inches(1.5), Inches(4.5), Inches(0.5),
         "Coverage Over Time", size=20, bold=True)
add_bullets(slide, Inches(7.5), Inches(2.2), Inches(4.5), Inches(4), [
    ("Week 1:", " most questions are new (\"I don't have that yet\")"),
    ("Month 3:", " 80% of questions get instant answers"),
    ("Month 12:", " 95% instant answers (comprehensive knowledge base)"),
    ("Steward workload", " decreases over time as coverage grows"),
], size=16)

# -- SLIDE 19: Section Divider --
slide = blank_slide()
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "DELIVERING VALUE", size=40, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "Self-Service + Dashboards, Together", size=22, color=GRAY, align=PP_ALIGN.CENTER)

# -- SLIDE 20: Delivery --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Self-Service + Dashboards: Not Either/Or", size=30, bold=True)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
         "The Data Agent doesn't replace dashboards — it eliminates the wait for them", size=18, color=GRAY)
add_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.5),
         "Self-Service (Data Agent)", size=22, bold=True)
add_bullets(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(3), [
    "Ask any question, get an answer in seconds",
    "No IT ticket, no backlog, no waiting",
    "Each user sees only what they're authorized to see",
    "Every answer is grounded in certified logic",
    "Usage patterns tell you what people actually need",
], size=16)
add_text(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(0.5),
         "Dashboards (Power BI)", size=22, bold=True)
add_bullets(slide, Inches(7), Inches(2.7), Inches(5.5), Inches(3), [
    "Recurring KPIs and executive views",
    "Built from the same certified knowledge base",
    "Numbers always match the agent's answers",
    "Prioritized by what users actually ask for",
    "High-demand metrics get promoted to dashboards",
], size=16)
add_text(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.8),
         "When a metric gets asked 100+ times, that's your signal to build a dashboard.\n"
         "Usage drives prioritization. No more guessing which reports to build.",
         size=16, color=GRAY, align=PP_ALIGN.CENTER)

# -- SLIDE 21: Section Divider --
slide = blank_slide()
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "MEASURING THE IMPACT", size=40, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "ROI That Compounds", size=22, color=GRAY, align=PP_ALIGN.CENTER)

# -- SLIDE 22: ROI --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Measuring Impact: The Numbers", size=30, bold=True)
add_text(slide, Inches(0.5), Inches(1.5), Inches(3), Inches(0.5),
         "Knowledge Base Coverage", size=16, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(2.2), Inches(3), Inches(0.6),
         "50 -> 250+", size=28, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(3.0), Inches(3), Inches(0.6),
         "Certified metrics grow\nfrom actual demand", size=14, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(3.5), Inches(1.5), Inches(3), Inches(0.5),
         "Instant Answer Rate", size=16, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(3.5), Inches(2.2), Inches(3), Inches(0.6),
         "40% -> 95%", size=28, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(3.5), Inches(3.0), Inches(3), Inches(0.6),
         "More questions answered\non the spot", size=14, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(6.5), Inches(1.5), Inches(3), Inches(0.5),
         "Time to Answer", size=16, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(6.5), Inches(2.2), Inches(3), Inches(0.6),
         "Weeks -> Seconds", size=28, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(6.5), Inches(3.0), Inches(3), Inches(0.6),
         "No more waiting for IT\nto build a report", size=14, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(9.5), Inches(1.5), Inches(3), Inches(0.5),
         "Steward Efficiency", size=16, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(9.5), Inches(2.2), Inches(3), Inches(0.6),
         "10x", size=28, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(9.5), Inches(3.0), Inches(3), Inches(0.6),
         "Demand-driven reviews\nnot committee meetings", size=14, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
         "Every number improves with time — the flywheel compounds",
         size=16, color=GRAY, align=PP_ALIGN.CENTER)

# -- SLIDE 23: Before & After --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Measuring Impact: The Before & After", size=30, bold=True)
add_text(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
         "BEFORE", size=24, bold=True)
add_bullets(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4), [
    "\"Which LOS number is right?\"",
    "3 analysts, 3 different answers",
    "6-week backlog for a new report",
    "Governance = a committee that meets monthly",
    "New metrics require a project proposal",
    "Nobody knows what's been asked before",
], size=17)
add_text(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
         "AFTER", size=24, bold=True)
add_bullets(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(4), [
    "\"Here's the certified ER LOS — asked 347 times this quarter\"",
    "One metric, one definition, traceable to source",
    "Ad-hoc answers in seconds, dashboards for top metrics",
    "Governance happens automatically from usage",
    "New metrics certified in days via steward review",
    "Usage data shows exactly what the organization cares about",
], size=17)

# -- SLIDE 24: Lessons Learned --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Lessons Learned", size=30, bold=True)
add_bullets(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5.5), [
    ("Start with what you have.", " Your existing reports and queries already contain the business logic. Extract it at scale — don't start from scratch."),
    ("\"I don't know\" is a feature, not a failure.", " In healthcare, a wrong answer is worse than no answer. And \"I don't know\" triggers the process that fills the gap."),
    ("Let usage drive prioritization.", " Don't guess which metrics matter — let the question volume tell you. Build dashboards for what people actually ask about."),
    ("Make governance invisible.", " If people have to stop working to do governance, they won't. Make it a byproduct of asking questions."),
    ("Security enables trust.", " When surgeons can only see their own data, they trust the system and use it more — which makes it grow faster."),
    ("The flywheel compounds.", " Early investment pays exponential returns. Month 12 looks nothing like month 1."),
], size=17)

# -- SLIDE 25: Key Takeaways --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Key Takeaways", size=30, bold=True)
add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.5),
         "1. Your existing reports are an untapped asset", size=22, bold=True)
add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(0.4),
         "The business logic is already there — extract it at scale to seed your certified knowledge base", size=17, color=GRAY)
add_text(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.5),
         "2. Usage is the best governance signal you have", size=22, bold=True)
add_text(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(0.4),
         "Every question reinforces a known metric or surfaces a new one — both make the system better", size=17, color=GRAY)
add_text(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(0.5),
         "3. Self-service and governance can be the same thing", size=22, bold=True)
add_text(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.4),
         "When asking a question IS the governance process, everyone participates without even knowing it", size=17, color=GRAY)
add_text(slide, Inches(1.5), Inches(5.4), Inches(10), Inches(0.5),
         "4. The full circle compounds over time", size=22, bold=True)
add_text(slide, Inches(1.5), Inches(5.9), Inches(10), Inches(0.4),
         "From FCOTS to a full OR efficiency suite — demand-driven growth that accelerates itself", size=17, color=GRAY)

# -- SLIDE 26: Thank You --
slide = blank_slide()
add_text(slide, Inches(0.8), Inches(2.0), Inches(10), Inches(1),
         "Thank You", size=48, bold=True)
add_text(slide, Inches(0.8), Inches(3.2), Inches(10), Inches(0.6),
         "Questions & Discussion", size=30, color=GRAY)
add_text(slide, Inches(0.8), Inches(4.5), Inches(10), Inches(0.5),
         "Sunny Zheng", size=22, bold=True)
add_text(slide, Inches(0.8), Inches(5.0), Inches(10), Inches(0.4),
         "HDAA Annual Conference  |  November 2026", size=16, color=GRAY)


# Speaker notes are in presentation/hdaa_speaker_notes.md
# (python-pptx notes_slide causes compatibility issues with PowerPoint for Mac)


# -- Save --
output_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "presentation")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "hdaa_full_circle_minimal.pptx")
prs.save(output_path)
print(f"Minimal deck saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
