"""
Generate HDAA Annual Conference Presentation Deck
"The Full Circle: From SQL Chaos to Self-Service Analytics"

Run: python3 scripts/generate_hdaa_deck.py
Output: presentation/hdaa_full_circle.pptx (also saves versioned copy)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import glob
import shutil

# -- Color palette --
DARK_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0x96, 0x88)
LIGHT_TEAL = RGBColor(0x4D, 0xB6, 0xAC)
ACCENT_ORANGE = RGBColor(0xFF, 0x7A, 0x2F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MEDIUM_GRAY = RGBColor(0x75, 0x75, 0x75)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=DARK_NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, left, top, width, height, items, font_size=18,
                     color=WHITE, bullet_color=TEAL):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.space_before = Pt(4)

        run_bullet = p.add_run()
        run_bullet.text = "\u25B8 "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Calibri"

        if isinstance(item, tuple):
            run_bold = p.add_run()
            run_bold.text = item[0]
            run_bold.font.size = Pt(font_size)
            run_bold.font.color.rgb = color
            run_bold.font.bold = True
            run_bold.font.name = "Calibri"

            run_rest = p.add_run()
            run_rest.text = item[1]
            run_rest.font.size = Pt(font_size)
            run_rest.font.color.rgb = color
            run_rest.font.name = "Calibri"
        else:
            run_text = p.add_run()
            run_text.text = item
            run_text.font.size = Pt(font_size)
            run_text.font.color.rgb = color
            run_text.font.name = "Calibri"
    return txBox


def section_divider(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_NAVY)
    add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Pt(4), fill_color=TEAL)
    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                title, font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.8),
                    subtitle, font_size=22, color=LIGHT_TEAL, alignment=PP_ALIGN.CENTER)
    return slide


def content_slide(title, bg_color=DARK_NAVY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, bg_color)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill_color=TEAL)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.8),
                title, font_size=30, color=WHITE, bold=True)
    add_shape(slide, Inches(0), Inches(1.1), SLIDE_W, Pt(3), fill_color=ACCENT_ORANGE)
    return slide


# ============================================================
# SLIDE 1 — Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-1), Inches(7), Inches(7))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x22, 0x3A, 0x5E)
circle.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(-0.3), Inches(5.6), Inches(5.6))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x28, 0x45, 0x6C)
circle2.line.fill.background()

add_shape(slide, Inches(0.8), Inches(4.4), Inches(4), Pt(4), fill_color=TEAL)

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(8), Inches(1.0),
            "The Full Circle:", font_size=28, color=LIGHT_TEAL, bold=False)
add_textbox(slide, Inches(0.8), Inches(2.5), Inches(10), Inches(1.5),
            "From SQL Chaos to\nSelf-Service AI Agent", font_size=48, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(4.7), Inches(8), Inches(0.5),
            "How Every User Question Builds Your Certified Data Foundation", font_size=20, color=LIGHT_TEAL)
add_textbox(slide, Inches(0.8), Inches(5.6), Inches(8), Inches(0.5),
            "Sunny Zheng", font_size=22, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(6.1), Inches(8), Inches(0.5),
            "HDAA Annual Conference  |  November 2026", font_size=16, color=MEDIUM_GRAY)


# ============================================================
# SLIDE 2 — About the Speaker
# ============================================================
slide = content_slide("About the Speaker")

add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), [
    ("BI Management Leader", " in Healthcare"),
    ("15+ years", " turning data into decisions for clinical and operational teams"),
    ("Microsoft Fabric ecosystem", " — Power BI, Data Agent, Notebooks"),
    ("Focused on", " making analytics self-service without sacrificing accuracy"),
    ("Passionate about", " making governance disappear into the background"),
], font_size=20)

box = add_rounded_rect(slide, Inches(7.5), Inches(2.0), Inches(4.5), Inches(3.5), fill_color=RGBColor(0x22, 0x3A, 0x5E))
add_textbox(slide, Inches(7.8), Inches(2.3), Inches(4), Inches(0.6),
            "The Challenge", font_size=24, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(7.8), Inches(3.0), Inches(4), Inches(2.2),
            "1,000+ reports and queries\nBusiness logic buried in code\nUsers waiting weeks for answers\nGovernance slowing everyone down",
            font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3 — Agenda
# ============================================================
slide = content_slide("Agenda")

agenda_items = [
    ("1.", " The Problem: Why Governance Feels Like a Tax"),
    ("2.", " The Idea: What If Usage Drove Governance?"),
    ("3.", " Building the Foundation: From Existing SQL to a Certified Knowledge Base"),
    ("4.", " The Flywheel: Ask, Answer, Certify, Grow"),
    ("5.", " Case Study: First Case On Time Start (FCOTS)"),
    ("6.", " The Growing Knowledge Base: Governance as a Byproduct"),
    ("7.", " Delivering Value: Self-Service + Dashboards"),
    ("8.", " Measuring the Impact: ROI That Compounds"),
    ("9.", " Lessons Learned & Q&A"),
]
add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5.5),
                 agenda_items, font_size=20, bullet_color=ACCENT_ORANGE)


# ============================================================
# SLIDE 4 — The Problem
# ============================================================
slide = content_slide("The Problem: Why Governance Feels Like a Tax")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "Sound familiar?", font_size=22, color=LIGHT_TEAL, bold=True)
add_bullet_slide(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), [
    "\"Can you pull me the readmission rate?\" — 6-week backlog",
    "Three analysts run the same metric, get three answers",
    "Governance team creates definitions nobody uses",
    "The more you govern, the slower everyone moves",
    "So governance falls behind, definitions go stale, trust erodes",
    "Result: governance and getting work done feel like opposites",
], font_size=18)

box = add_rounded_rect(slide, Inches(7.5), Inches(1.8), Inches(5), Inches(4.5), fill_color=RGBColor(0x22, 0x3A, 0x5E))
add_textbox(slide, Inches(7.8), Inches(2.0), Inches(4.4), Inches(0.5),
            "The Vicious Cycle", font_size=24, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

cycle_items = [
    ("Governance is slow", ACCENT_ORANGE),
    ("People work around it", ACCENT_ORANGE),
    ("Definitions go stale", ACCENT_ORANGE),
    ("Trust in data declines", ACCENT_ORANGE),
    ("More governance imposed", ACCENT_ORANGE),
]
for i, (label, color) in enumerate(cycle_items):
    y = Inches(2.8) + Inches(0.55) * i
    indicator = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.0), y, Inches(0.3), Inches(0.3))
    indicator.fill.solid()
    indicator.fill.fore_color.rgb = color
    indicator.line.fill.background()
    add_textbox(slide, Inches(8.5), y, Inches(3.5), Inches(0.35),
                label, font_size=17, color=WHITE)

add_textbox(slide, Inches(8.0), Inches(5.7), Inches(4.0), Inches(0.4),
            "(repeat forever)", font_size=15, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5 — The Big Idea
# ============================================================
slide = content_slide("The Idea: What If Usage Drove Governance?")

add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
            "What if every question a user asks\nmakes your data foundation stronger?",
            font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(3), Inches(3.3), Inches(7), Pt(3), fill_color=TEAL)

add_textbox(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(2.5),
            "Instead of governance being a gate people pass through,\n"
            "it becomes a knowledge base that grows from demand.\n\n"
            "Every question either reinforces what we know\n"
            "or surfaces what we don't — and kicks off a review.",
            font_size=22, color=LIGHT_TEAL, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(6.2), Inches(9), Inches(0.6),
            "Governance is no longer a chore. It's a byproduct of people doing their jobs.",
            font_size=20, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6 — Section Divider: BUILD THE FOUNDATION
# ============================================================
section_divider("BUILDING THE FOUNDATION", "From Existing SQL to a Certified Knowledge Base")


# ============================================================
# SLIDE 7 — Seed: Extract
# ============================================================
slide = content_slide("The Foundation: Your SQL Already Contains the Answers")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
            "You don't start from scratch — your existing queries are the raw material",
            font_size=20, color=LIGHT_TEAL)

add_bullet_slide(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), [
    ("Scan your existing reports and queries", " — automatically extract the business logic"),
    ("AI drafts plain-English definitions", " from the logic it finds"),
    ("Your team reviews in bulk", " — not one definition at a time"),
    ("Stewards certify", " the business meaning: approve, modify, or reject"),
    ("Result:", " an initial certified knowledge base of your most-used metrics"),
    ("This is the seed.", " What makes it grow is what happens next."),
], font_size=18)

code_box = add_rounded_rect(slide, Inches(7.2), Inches(2.0), Inches(5.3), Inches(4.2),
                            fill_color=RGBColor(0x1A, 0x1A, 0x2E))
add_textbox(slide, Inches(7.5), Inches(2.1), Inches(5), Inches(0.4),
            "Example: One report yields...", font_size=14, color=MEDIUM_GRAY)
add_textbox(slide, Inches(7.5), Inches(2.6), Inches(5), Inches(3.5),
            "Report: OR Efficiency Dashboard\n"
            "\n"
            "  3 business metrics found:\n"
            "    \"First Case On Time Start\"\n"
            "    \"Turnover Time\"\n"
            "    \"Case Duration Accuracy\"\n"
            "\n"
            "  Each linked to its source\n"
            "  tables and calculation logic\n"
            "\n"
            "  -> Sent to stewards for review",
            font_size=16, color=RGBColor(0xA0, 0xE0, 0xA0), font_name="Consolas")


# ============================================================
# SLIDE 8 — How It's Organized
# ============================================================
slide = content_slide("The Foundation: How the Knowledge Base Is Organized")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
            "Three layers connect business questions to source data — with full traceability",
            font_size=20, color=LIGHT_TEAL)

layers = [
    ("BUSINESS METRICS", "Certified definitions (e.g., First Case On Time Start)\nOwned by a steward  |  Weighted by how often people ask",
     Inches(2.5), Inches(2.3), Inches(8), Inches(1.2), TEAL),
    ("CALCULATION LOGIC", "The step-by-step logic behind each metric\nSmall, reusable pieces — not giant SQL blocks",
     Inches(2.5), Inches(4.0), Inches(8), Inches(1.2), RGBColor(0x26, 0xA6, 0x9A)),
    ("SOURCE DATA", "The actual tables, columns, and filters\nEnriched with plain-English descriptions",
     Inches(2.5), Inches(5.7), Inches(8), Inches(1.2), RGBColor(0x00, 0x60, 0x56)),
]

for label, desc, left, top, width, height, color in layers:
    box = add_rounded_rect(slide, left, top, width, height, fill_color=color)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.1), width - Inches(0.6), Inches(0.4),
                label, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.5), width - Inches(0.6), Inches(0.6),
                desc, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

for y_top in [Inches(3.5), Inches(5.2)]:
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.3), y_top, Inches(0.4), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT_ORANGE
    arrow.line.fill.background()


# ============================================================
# SLIDE 9 — Section Divider: THE FLYWHEEL
# ============================================================
section_divider("THE FLYWHEEL", "Ask, Answer, Certify, Grow")


# ============================================================
# SLIDE 10 — The Flywheel Overview
# ============================================================
slide = content_slide("The Flywheel: Every Question Makes the System Smarter")

flywheel = [
    ("USER ASKS\nA QUESTION", Inches(5.2), Inches(1.5), TEAL),
    ("AGENT\nANSWERS", Inches(9.0), Inches(3.2), LIGHT_TEAL),
    ("KNOWLEDGE\nBASE GROWS", Inches(5.2), Inches(5.0), RGBColor(0x26, 0xA6, 0x9A)),
    ("BETTER\nANSWERS", Inches(1.5), Inches(3.2), ACCENT_ORANGE),
]

for label, left, top, color in flywheel:
    box = add_rounded_rect(slide, left, top, Inches(2.5), Inches(1.3), fill_color=color)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.15), Inches(2.3), Inches(1.0),
                label, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.7), Inches(2.1), Inches(1.2), Inches(0.4))
arrow.fill.solid()
arrow.fill.fore_color.rgb = WHITE
arrow.line.fill.background()

arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.9), Inches(4.5), Inches(0.4), Inches(0.6))
arrow.fill.solid()
arrow.fill.fore_color.rgb = WHITE
arrow.line.fill.background()

arrow = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(3.9), Inches(5.4), Inches(1.2), Inches(0.4))
arrow.fill.solid()
arrow.fill.fore_color.rgb = WHITE
arrow.line.fill.background()

arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(2.4), Inches(2.0), Inches(0.4), Inches(0.6))
arrow.fill.solid()
arrow.fill.fore_color.rgb = WHITE
arrow.line.fill.background()

add_textbox(slide, Inches(4.0), Inches(3.2), Inches(4.5), Inches(1.2),
            "FULL\nCIRCLE", font_size=36, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
            "No governance committees. No spreadsheets. The knowledge base grows because people use it.",
            font_size=18, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11 — How It Works: The Two Paths
# ============================================================
slide = content_slide("Two Paths: Every Question Has a Destination")

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.6),
            "When someone asks a question, one of two things happens:",
            font_size=22, color=LIGHT_TEAL, alignment=PP_ALIGN.CENTER)

# Path A box
path_a = add_rounded_rect(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4.2),
                          fill_color=RGBColor(0x22, 0x3A, 0x5E))
add_textbox(slide, Inches(1.0), Inches(2.7), Inches(5.1), Inches(0.5),
            "PATH A: We Already Know This", font_size=24, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(1.3), Inches(3.3), Inches(4.5), Pt(2), fill_color=TEAL)
add_bullet_slide(slide, Inches(1.3), Inches(3.6), Inches(4.8), Inches(3.0), [
    "The answer exists in the certified knowledge base",
    "User gets an instant, accurate answer",
    "The metric gains weight — it's clearly important",
    "Most-asked metrics become visible to leadership",
    "High-demand metrics get promoted to dashboards",
], font_size=16, bullet_color=TEAL)

# Path B box
path_b = add_rounded_rect(slide, Inches(7.0), Inches(2.5), Inches(5.5), Inches(4.2),
                          fill_color=RGBColor(0x22, 0x3A, 0x5E))
add_textbox(slide, Inches(7.2), Inches(2.7), Inches(5.1), Inches(0.5),
            "PATH B: We Don't Know This Yet", font_size=24, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(7.5), Inches(3.3), Inches(4.5), Pt(2), fill_color=ACCENT_ORANGE)
add_bullet_slide(slide, Inches(7.5), Inches(3.6), Inches(4.8), Inches(3.0), [
    "The agent says \"I don't have that yet\"",
    "A request is sent to the data steward",
    "Steward reviews and certifies the definition",
    "New metric added to the knowledge base",
    "Next time anyone asks, they get an instant answer",
], font_size=16, bullet_color=ACCENT_ORANGE)

add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
            "Both paths make the system better. There is no wasted question.",
            font_size=20, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12 — Section Divider: CASE STUDY
# ============================================================
section_divider("CASE STUDY", "First Case On Time Start (FCOTS)")


# ============================================================
# SLIDE 13 — FCOTS: The Story
# ============================================================
slide = content_slide("Case Study: First Case On Time Start")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
            "How one certified metric helped transform surgical efficiency",
            font_size=20, color=LIGHT_TEAL)

# Left side - the story
add_bullet_slide(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), [
    ("FCOTS =", " % of first surgical cases that start on time"),
    ("3 years ago:", " our rate was just over 20%"),
    ("Today:", " we've reached 60%+"),
    ("Goal:", " 90%+ on-time start rate"),
    ("Working with", " the Surgery Medical Director to drive improvement"),
    ("One certified definition,", " one source of truth, consistent tracking over time"),
], font_size=18)

# Right side - progress visualization
box = add_rounded_rect(slide, Inches(7.5), Inches(1.8), Inches(5), Inches(5),
                       fill_color=RGBColor(0x22, 0x3A, 0x5E))
add_textbox(slide, Inches(7.8), Inches(2.0), Inches(4.4), Inches(0.5),
            "FCOTS Journey", font_size=20, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# Progress bars
milestones = [
    ("3 Years Ago", "20%+", Inches(1.0), ACCENT_ORANGE),
    ("Today", "60%+", Inches(2.8), TEAL),
    ("Goal", "90%+", Inches(4.2), RGBColor(0x4C, 0xAF, 0x50)),
]
for i, (label, pct, bar_width, color) in enumerate(milestones):
    y = Inches(2.8) + Inches(0.9) * i
    add_textbox(slide, Inches(7.8), y, Inches(1.5), Inches(0.4),
                label, font_size=14, color=MEDIUM_GRAY, bold=True)
    bar = add_rounded_rect(slide, Inches(7.8), y + Inches(0.35), bar_width, Inches(0.35), fill_color=color)
    add_textbox(slide, Inches(7.8) + bar_width + Inches(0.15), y + Inches(0.35), Inches(1.0), Inches(0.35),
                pct, font_size=16, color=color, bold=True)

add_textbox(slide, Inches(7.8), Inches(5.8), Inches(4.4), Inches(0.8),
            "This improvement was only possible\nbecause everyone trusted the same number.",
            font_size=15, color=LIGHT_TEAL, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 14 — FCOTS: The Challenge & the Flywheel Solution
# ============================================================
slide = content_slide("FCOTS: Why This Metric Needs Self-Service")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
            "The Medical Director's challenge: individual performance is sensitive data",
            font_size=20, color=LIGHT_TEAL)

# The challenge
challenge_box = add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(2.2),
                                 fill_color=RGBColor(0x22, 0x3A, 0x5E))
add_textbox(slide, Inches(1.0), Inches(2.5), Inches(5.1), Inches(0.5),
            "The Challenge", font_size=20, color=ACCENT_ORANGE, bold=True)
add_bullet_slide(slide, Inches(1.0), Inches(3.1), Inches(5.1), Inches(1.2), [
    "Medical Director needs 1:1s with surgeons about their FCOTS",
    "Can't share a dashboard showing everyone's performance",
    "Individual surgeon data is sensitive — requires privacy",
], font_size=16, bullet_color=ACCENT_ORANGE)

# The solution
solution_box = add_rounded_rect(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.2),
                                fill_color=RGBColor(0x22, 0x3A, 0x5E))
add_textbox(slide, Inches(1.0), Inches(5.0), Inches(5.1), Inches(0.5),
            "The Self-Service Solution", font_size=20, color=TEAL, bold=True)
add_bullet_slide(slide, Inches(1.0), Inches(5.6), Inches(5.1), Inches(1.2), [
    "Surgeons ask the Data Agent: \"What's my FCOTS rate?\"",
    "Each surgeon only sees their own data — built-in security",
    "Medical Director sees the department-level view",
], font_size=16, bullet_color=TEAL)

# Right side - three views
box = add_rounded_rect(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(5.2),
                       fill_color=RGBColor(0x22, 0x3A, 0x5E))
add_textbox(slide, Inches(7.3), Inches(2.2), Inches(4.9), Inches(0.5),
            "Same Metric, Right Access", font_size=20, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

# Three user views
views = [
    ("Dr. Smith asks:", "\"What's my FCOTS this month?\"", "Answer: 78% (your cases only)", TEAL),
    ("Dr. Jones asks:", "\"What's my FCOTS this month?\"", "Answer: 62% (your cases only)", LIGHT_TEAL),
    ("Medical Director asks:", "\"Show me department FCOTS\"", "Answer: 68% overall + trends", ACCENT_ORANGE),
]
for i, (who, question, answer, color) in enumerate(views):
    y = Inches(3.0) + Inches(1.3) * i
    add_textbox(slide, Inches(7.3), y, Inches(4.9), Inches(0.35),
                who, font_size=14, color=color, bold=True)
    add_textbox(slide, Inches(7.3), y + Inches(0.35), Inches(4.9), Inches(0.35),
                question, font_size=13, color=WHITE)
    add_textbox(slide, Inches(7.3), y + Inches(0.7), Inches(4.9), Inches(0.35),
                answer, font_size=13, color=RGBColor(0xA0, 0xE0, 0xA0))

add_textbox(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.4),
            "One certified definition. Personalized access. No IT tickets. No privacy concerns.",
            font_size=17, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 15 — FCOTS: The Full Circle in Action
# ============================================================
slide = content_slide("FCOTS: The Full Circle in Action")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
            "What happens when surgeons start using the agent regularly:",
            font_size=20, color=LIGHT_TEAL)

add_bullet_slide(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), [
    ("Surgeons ask about their FCOTS", " — the metric gains weight in the system"),
    ("Some ask new questions:", " \"What's my average delay time?\" \"Which rooms start late most?\""),
    ("Those questions trigger", " steward review — new metrics get certified"),
    ("The knowledge base grows", " around OR efficiency — driven by surgeon demand"),
    ("Medical Director sees", " which metrics surgeons care about most"),
    ("Result:", " a richer, more useful system built from actual clinical needs"),
], font_size=18)

# Right side - the mini flywheel for FCOTS
box = add_rounded_rect(slide, Inches(7.5), Inches(1.8), Inches(5), Inches(5),
                       fill_color=RGBColor(0x22, 0x3A, 0x5E))
add_textbox(slide, Inches(7.8), Inches(2.0), Inches(4.4), Inches(0.5),
            "FCOTS Flywheel Growth", font_size=20, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

growth = [
    ("Start:", "1 certified metric\n(FCOTS)", MEDIUM_GRAY),
    ("Month 1:", "Surgeons start asking\nabout delay reasons", TEAL),
    ("Month 3:", "5 related OR efficiency\nmetrics certified", RGBColor(0x26, 0xA6, 0x9A)),
    ("Month 6:", "Full OR efficiency suite\nbuilt from surgeon demand", RGBColor(0x4C, 0xAF, 0x50)),
]
for i, (when, what, color) in enumerate(growth):
    y = Inches(2.8) + Inches(1.0) * i
    indicator = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.0), y + Inches(0.05), Inches(0.3), Inches(0.3))
    indicator.fill.solid()
    indicator.fill.fore_color.rgb = color
    indicator.line.fill.background()
    add_textbox(slide, Inches(8.5), y, Inches(3.7), Inches(0.3),
                when, font_size=14, color=color, bold=True)
    add_textbox(slide, Inches(8.5), y + Inches(0.3), Inches(3.7), Inches(0.5),
                what, font_size=13, color=WHITE)


# ============================================================
# SLIDE 16 — Section Divider: GROWING KNOWLEDGE BASE
# ============================================================
section_divider("THE GROWING KNOWLEDGE BASE", "Governance as a Byproduct of Work")


# ============================================================
# SLIDE 17 — The Growing Knowledge Base
# ============================================================
slide = content_slide("The Knowledge Base Grows From Demand, Not Committees")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "Month over month, your certified metrics expand:", font_size=20, color=LIGHT_TEAL, bold=True)

months = [
    ("Month 1", "Seed from existing\nreports & queries", "50", MEDIUM_GRAY),
    ("Month 3", "User questions add\n30 new certified metrics", "80", TEAL),
    ("Month 6", "Steward queue drives\n60 more certifications", "140", RGBColor(0x26, 0xA6, 0x9A)),
    ("Month 12", "Flywheel matures\n250+ certified metrics", "250+", RGBColor(0x4C, 0xAF, 0x50)),
]

for i, (month, desc, count, color) in enumerate(months):
    left = Inches(0.8) + Inches(3.0) * i

    bar_heights = [Inches(1.0), Inches(1.6), Inches(2.8), Inches(4.0)]
    bar_top = Inches(6.2) - bar_heights[i]
    bar = add_rounded_rect(slide, left + Inches(0.3), bar_top, Inches(2.0), bar_heights[i], fill_color=color)

    add_textbox(slide, left + Inches(0.3), bar_top + Inches(0.1), Inches(2.0), Inches(0.5),
                count, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, left, Inches(2.2), Inches(2.6), Inches(0.3),
                month, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
            "Traditional governance: ~50 definitions/year.  Demand-driven: 250+ and accelerating.",
            font_size=18, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 18 — Why It Accelerates
# ============================================================
slide = content_slide("Why the Flywheel Accelerates Over Time")

add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), [
    ("More certified metrics", " = more questions the agent can answer"),
    ("More answers", " = more people trust and adopt the agent"),
    ("More users", " = more questions — both known and new"),
    ("More new questions", " = more steward certifications"),
    ("More certifications", " = even more metrics the agent knows"),
    ("Each cycle is faster.", " The foundation compounds."),
], font_size=18)

box = add_rounded_rect(slide, Inches(7.5), Inches(1.8), Inches(5), Inches(4.8),
                       fill_color=RGBColor(0x22, 0x3A, 0x5E))
add_textbox(slide, Inches(7.8), Inches(2.0), Inches(4.4), Inches(0.5),
            "Coverage Over Time", font_size=20, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

add_bullet_slide(slide, Inches(7.8), Inches(2.8), Inches(4.4), Inches(3.5), [
    ("Week 1:", " most questions are new\n(\"I don't have that yet\")"),
    ("Month 3:", " 80% of questions get\ninstant answers"),
    ("Month 12:", " 95% instant answers\n(comprehensive knowledge base)"),
    ("Steward workload", " decreases over time\nas coverage grows"),
], font_size=16, bullet_color=TEAL)


# ============================================================
# SLIDE 19 — Section Divider: DELIVERING VALUE
# ============================================================
section_divider("DELIVERING VALUE", "Self-Service + Dashboards, Together")


# ============================================================
# SLIDE 20 — Delivery: Agent + Reports
# ============================================================
slide = content_slide("Self-Service + Dashboards: Not Either/Or")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
            "The Data Agent doesn't replace dashboards — it eliminates the wait for them",
            font_size=20, color=LIGHT_TEAL)

agent_box = add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(3.8),
                             fill_color=RGBColor(0x22, 0x3A, 0x5E))
add_textbox(slide, Inches(1.0), Inches(2.5), Inches(5.1), Inches(0.5),
            "Self-Service (Data Agent)", font_size=22, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_slide(slide, Inches(1.3), Inches(3.2), Inches(4.8), Inches(2.8), [
    "Ask any question, get an answer in seconds",
    "No IT ticket, no backlog, no waiting",
    "Each user sees only what they're authorized to see",
    "Every answer is grounded in certified logic",
    "Usage patterns tell you what people actually need",
], font_size=16, bullet_color=TEAL)

report_box = add_rounded_rect(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(3.8),
                              fill_color=RGBColor(0x22, 0x3A, 0x5E))
add_textbox(slide, Inches(7.2), Inches(2.5), Inches(5.1), Inches(0.5),
            "Dashboards (Power BI)", font_size=22, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_slide(slide, Inches(7.5), Inches(3.2), Inches(4.8), Inches(2.8), [
    "Recurring KPIs and executive views",
    "Built from the same certified knowledge base",
    "Numbers always match the agent's answers",
    "Prioritized by what users actually ask for",
    "High-demand metrics get promoted to dashboards",
], font_size=16, bullet_color=ACCENT_ORANGE)

add_textbox(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.8),
            "When a metric gets asked 100+ times, that's your signal to build a dashboard.\n"
            "Usage drives prioritization. No more guessing which reports to build.",
            font_size=17, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 21 — Section Divider: MEASURING
# ============================================================
section_divider("MEASURING THE IMPACT", "ROI That Compounds")


# ============================================================
# SLIDE 22 — Measure: Quantitative ROI
# ============================================================
slide = content_slide("Measuring Impact: The Numbers")

kpis = [
    ("Knowledge Base\nCoverage", "50 -> 250+", "Certified metrics\ngrow from actual demand"),
    ("Instant Answer\nRate", "40% -> 95%", "More questions\nanswered on the spot"),
    ("Time to Answer", "Weeks -> Seconds", "No more waiting\nfor IT to build a report"),
    ("Steward\nEfficiency", "10x", "Demand-driven reviews\nnot committee meetings"),
]

for i, (label, value, desc) in enumerate(kpis):
    col = i % 4
    left = Inches(0.5) + Inches(3.1) * col
    top = Inches(1.8)

    box = add_rounded_rect(slide, left, top, Inches(2.8), Inches(3.5),
                           fill_color=RGBColor(0x22, 0x3A, 0x5E))
    add_textbox(slide, left + Inches(0.1), top + Inches(0.2), Inches(2.6), Inches(0.8),
                label, font_size=16, color=LIGHT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(1.1), Inches(2.6), Inches(0.8),
                value, font_size=32, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(2.2), Inches(2.6), Inches(1.0),
                desc, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
            "Every number improves with time — the flywheel compounds",
            font_size=18, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 23 — Measure: Before & After
# ============================================================
slide = content_slide("Measuring Impact: The Before & After")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "BEFORE", font_size=26, color=ACCENT_ORANGE, bold=True)
add_shape(slide, Inches(0.8), Inches(2.1), Inches(5.2), Pt(2), fill_color=ACCENT_ORANGE)
add_bullet_slide(slide, Inches(0.8), Inches(2.3), Inches(5.2), Inches(4.5), [
    "\"Which LOS number is right?\"",
    "3 analysts, 3 different answers",
    "6-week backlog for a new report",
    "Governance = a committee that meets monthly",
    "New metrics require a project proposal",
    "Nobody knows what's been asked before",
], font_size=17, bullet_color=ACCENT_ORANGE)

add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
            "AFTER", font_size=26, color=TEAL, bold=True)
add_shape(slide, Inches(7.0), Inches(2.1), Inches(5.2), Pt(2), fill_color=TEAL)
add_bullet_slide(slide, Inches(7.0), Inches(2.3), Inches(5.2), Inches(4.5), [
    "\"Here's the certified ER LOS — asked 347 times this quarter\"",
    "One metric, one definition, traceable to source",
    "Ad-hoc answers in seconds, dashboards for top metrics",
    "Governance happens automatically from usage",
    "New metrics certified in days via steward review",
    "Usage data shows exactly what the organization cares about",
], font_size=17, bullet_color=TEAL)


# ============================================================
# SLIDE 24 — Lessons Learned
# ============================================================
slide = content_slide("Lessons Learned")

add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5.5), [
    ("Start with what you have.", " Your existing reports and queries already contain the business logic. "
     "Extract it at scale — don't start from scratch."),
    ("\"I don't know\" is a feature, not a failure.", " In healthcare, a wrong answer is worse than "
     "no answer. And \"I don't know\" triggers the process that fills the gap."),
    ("Let usage drive prioritization.", " Don't guess which metrics matter — let the question volume "
     "tell you. Build dashboards for what people actually ask about."),
    ("Make governance invisible.", " If people have to stop working to do governance, they won't. "
     "Make it a byproduct of asking questions."),
    ("Security enables trust.", " When surgeons can only see their own data, "
     "they trust the system and use it more — which makes it grow faster."),
    ("The flywheel compounds.", " Early investment pays exponential returns. "
     "Month 12 looks nothing like month 1."),
], font_size=18)


# ============================================================
# SLIDE 25 — Key Takeaways
# ============================================================
slide = content_slide("Key Takeaways")

takeaways = [
    ("1", "Your existing reports are an untapped asset",
     "The business logic is already there — extract it at scale to seed your certified knowledge base"),
    ("2", "Usage is the best governance signal you have",
     "Every question reinforces a known metric or surfaces a new one — both make the system better"),
    ("3", "Self-service and governance can be the same thing",
     "When asking a question IS the governance process, everyone participates without even knowing it"),
    ("4", "The full circle compounds over time",
     "From FCOTS to a full OR efficiency suite — demand-driven growth that accelerates itself"),
]

for i, (num, title, desc) in enumerate(takeaways):
    y = Inches(1.6) + Inches(1.35) * i
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.1), Inches(0.6), Inches(0.6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TEAL
    circ.line.fill.background()
    add_textbox(slide, Inches(0.8), y + Inches(0.12), Inches(0.6), Inches(0.6),
                num, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.7), y, Inches(10.5), Inches(0.5),
                title, font_size=22, color=WHITE, bold=True)
    add_textbox(slide, Inches(1.7), y + Inches(0.5), Inches(10.5), Inches(0.5),
                desc, font_size=17, color=LIGHT_TEAL)


# ============================================================
# SLIDE 26 — Q&A / Contact
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(1), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x22, 0x3A, 0x5E)
circle.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(1.7), Inches(4.6), Inches(4.6))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x28, 0x45, 0x6C)
circle2.line.fill.background()

add_shape(slide, Inches(0.8), Inches(4.4), Inches(4), Pt(4), fill_color=TEAL)

add_textbox(slide, Inches(0.8), Inches(2.0), Inches(8), Inches(1.0),
            "Thank You", font_size=48, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(3.2), Inches(8), Inches(1.0),
            "Questions & Discussion", font_size=30, color=LIGHT_TEAL)
add_textbox(slide, Inches(0.8), Inches(4.8), Inches(8), Inches(0.5),
            "Sunny Zheng", font_size=22, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(5.4), Inches(8), Inches(0.5),
            "HDAA Annual Conference  |  November 2026", font_size=16, color=MEDIUM_GRAY)


# ============================================================
# Save with auto-versioning
# ============================================================
output_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "presentation")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "hdaa_full_circle.pptx")

# Auto-version: if the file already exists, save a versioned backup first
if os.path.exists(output_path):
    existing_versions = glob.glob(os.path.join(output_dir, "hdaa_full_circle_v*.pptx"))
    version_nums = []
    for v in existing_versions:
        try:
            num = int(os.path.basename(v).replace("hdaa_full_circle_v", "").replace(".pptx", ""))
            version_nums.append(num)
        except ValueError:
            pass
    next_version = max(version_nums, default=0) + 1
    backup_path = os.path.join(output_dir, f"hdaa_full_circle_v{next_version}.pptx")
    shutil.copy2(output_path, backup_path)
    print(f"Previous version saved as: {backup_path}")

prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
